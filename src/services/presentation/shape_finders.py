# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Assistance Micro Design
"""Fonctions de detection de shapes dans les slides PowerPoint.

Fonctions partagees entre PresentationEditor et FileInspector
pour garantir une detection coherente des titres, sous-titres et puces.
"""

from __future__ import annotations

from typing import TYPE_CHECKING, Any


if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide


def find_title_shape(slide: Slide) -> BaseShape | None:
    """Trouve le premier shape contenant du texte (titre presume).

    Args:
        slide: Slide PowerPoint a analyser.

    Returns:
        Le shape titre ou None si aucun shape texte.
    """
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text:
            return shape
    return None


def find_subtitle_shape(slide: Slide) -> BaseShape | None:
    """Trouve le deuxieme shape contenant du texte (sous-titre presume).

    Args:
        slide: Slide PowerPoint a analyser.

    Returns:
        Le shape sous-titre ou None si moins de 2 shapes texte.
    """
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text]
    if len(text_shapes) >= 2:
        return text_shapes[1]
    return None


def find_bullets_shape(slide: Slide) -> BaseShape | None:
    """Trouve le shape contenant des puces (le plus grand textbox apres le titre).

    Args:
        slide: Slide PowerPoint a analyser.

    Returns:
        Le shape puces ou None si moins de 2 shapes texte.
    """
    text_shapes: list[Any] = [s for s in slide.shapes if s.has_text_frame]
    if len(text_shapes) >= 2:
        return text_shapes[1]
    return None
