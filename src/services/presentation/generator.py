# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Assistance Micro Design
"""Generateur de presentations PowerPoint (.pptx)."""

from __future__ import annotations

import asyncio
import logging
from copy import deepcopy
from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from src.core.config import settings
from src.core.exceptions import (
    PresentationChartError,
    PresentationGenerationError,
    PresentationImageNotFoundError,
    PresentationOutputTooLargeError,
    PresentationTemplateNotFoundError,
)
from src.models.api import CreatePresentationParams, CreatePresentationResult
from src.models.presentation_generation import (
    BulletItem,
    ChartSlideDef,
    ClosingSlideDef,
    ContentBulletsSlideDef,
    ContentWithImageSlideDef,
    ImageDef,
    ImageFullSlideDef,
    PresentationChartDef,
    SectionHeaderSlideDef,
    TextStyle,
    TitleSlideDef,
    TwoColumnsSlideDef,
)


if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.text.text import TextFrame

logger = logging.getLogger(__name__)


# Mapping noms simplifies -> enums python-pptx (XY_SCATTER utilise car SCATTER n'existe pas)
_CHART_TYPE_MAP: dict[str, Any] = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}

# Index du layout "Blank" dans la collection par defaut de python-pptx
_BLANK_LAYOUT_INDEX = 6

# Mapping: notre layout -> noms de layouts cherches dans le template (par priorite)
# Couvre les noms standard Office, Google Slides et variantes courantes
_LAYOUT_NAME_CANDIDATES: dict[str, list[str]] = {
    "title_slide": ["Title Slide", "TITLE", "Title"],
    "content_bullets": [
        "Title and Content",
        "OBJECT",
        "Title, Content",
        "Content",
    ],
    "section_header": ["Section Header", "SECTION_HEADER", "Section"],
    "two_columns": [
        "Two Content",
        "TWO_OBJECTS",
        "Two Objects",
        "Comparison",
    ],
    "content_with_image": [
        "Title and Content",
        "OBJECT",
        "Content with Caption",
        "OBJECT_WITH_CAPTION_TEXT",
    ],
    "chart_slide": ["Title and Content", "OBJECT", "Title Only", "TITLE_ONLY"],
    "closing": ["Title Slide", "TITLE", "Title"],
    "image_full": ["Blank", "BLANK"],
}

# Placeholder idx constants
_PH_TITLE = 0
_PH_BODY = 1
_PH_BODY_RIGHT = 2


class PresentationGenerator:
    """Genere des presentations PowerPoint a partir de definitions structurees.

    Utilise python-pptx pour creer des fichiers .pptx avec slides,
    texte, images, graphiques et notes.

    Attributes:
        _output_path: Repertoire de sortie des fichiers generes.
        _images_path: Repertoire des images disponibles.
        _templates_path: Repertoire des templates disponibles.
    """

    def __init__(
        self,
        output_path: Path | None = None,
        images_path: Path | None = None,
        templates_path: Path | None = None,
    ) -> None:
        """Initialise le generateur.

        Args:
            output_path: Repertoire de sortie (defaut: settings.OUTPUT_PATH).
            images_path: Repertoire des images (defaut: settings.IMAGES_POWERPOINT_PATH).
            templates_path: Repertoire des templates (defaut: settings.TEMPLATES_PPTX_PATH).
        """
        self._output_path = Path(output_path or settings.OUTPUT_PATH)
        self._images_path = Path(images_path or settings.IMAGES_POWERPOINT_PATH)
        self._templates_path = Path(templates_path or settings.TEMPLATES_PPTX_PATH)
        self._max_output_size_mb = settings.MAX_OUTPUT_FILE_SIZE_MB
        self._layout_index: dict[str, Any] = {}
        self._has_template = False
        self._reference_slides: dict[int, dict[str, Any]] = {}
        self._template_slide_map: list[int] | None = None
        self._current_slide_idx = 0
        self._last_slide_cloned = False

    async def generate(self, params: CreatePresentationParams) -> CreatePresentationResult:
        """Point d'entree principal. Cree le fichier pptx.

        Args:
            params: Parametres de creation de la presentation.

        Returns:
            Resultat avec chemin, stats et taille du fichier.
        """
        return await asyncio.to_thread(self._generate_sync, params)

    def _generate_sync(self, params: CreatePresentationParams) -> CreatePresentationResult:
        """Generation synchrone du fichier (appelee via to_thread).

        Returns:
            Resultat avec chemin, stats et taille du fichier.
        """
        self.ensure_output_dir()
        safe_filename = self.sanitize_filename(params.filename)
        file_path = self._output_path / safe_filename
        overwritten = file_path.exists()

        prs = self._load_or_create_presentation(params.template)
        self._layout_index = self._build_layout_index(prs)
        self._has_template = params.template is not None
        self._template_slide_map = params.template_slide_map
        self._current_slide_idx = 0
        total_images, total_charts = self._process_slides(prs, params.slides)

        if params.author:
            prs.core_properties.author = params.author

        file_size = self.save_and_verify(prs, file_path, safe_filename)

        logger.info(
            "Presentation generee: %s (%d slides, %d images, %d charts, %d octets)",
            safe_filename,
            len(params.slides),
            total_images,
            total_charts,
            file_size,
        )

        return CreatePresentationResult(
            file_path=str(file_path),
            filename=safe_filename,
            slides_created=len(params.slides),
            total_images=total_images,
            total_charts=total_charts,
            file_size_bytes=file_size,
            overwritten=overwritten,
        )

    def _load_or_create_presentation(self, template: str | None) -> Presentation:
        """Charge un template ou cree une presentation vierge.

        Args:
            template: Nom du fichier template (optionnel).

        Returns:
            Instance Presentation.

        Raises:
            PresentationTemplateNotFoundError: Si le template est introuvable.
        """
        if not template:
            self._reference_slides = {}
            return Presentation()

        if ".." in template or "/" in template or "\\" in template:
            raise PresentationGenerationError(
                message=f"Nom de template invalide: {template}",
                code="INVALID_TEMPLATE_NAME",
                suggestion="Le nom du template ne doit pas contenir '..' , '/' ou '\\'.",
                parameter="template",
                retry=True,
            )

        template_path = self._templates_path / template
        if not template_path.exists():
            available = self._list_templates()
            raise PresentationTemplateNotFoundError(template, available)
        prs = Presentation(str(template_path))
        self._build_reference_slides(prs)
        self._remove_existing_slides(prs)
        return prs

    def _list_templates(self) -> list[str]:
        """Liste les templates disponibles."""
        if not self._templates_path.exists():
            return []
        return sorted(f.name for f in self._templates_path.glob("*.pptx"))

    def _remove_existing_slides(self, prs: Presentation) -> None:
        """Supprime tous les slides existants du template.

        Conserve les layouts et le theme, mais retire les slides de contenu
        pour ne garder que la structure du template.

        Args:
            prs: Instance Presentation dont les slides seront supprimes.
        """
        slide_id_list = prs.slides._sldIdLst
        while len(slide_id_list):
            r_id = slide_id_list[-1].rId
            prs.part.drop_rel(r_id)
            del slide_id_list[-1]

    def _build_reference_slides(self, prs: Presentation) -> None:
        """Stocke les donnees XML et relations de chaque slide du template.

        Doit etre appele AVANT _remove_existing_slides pour capturer
        le design des slides avant suppression.

        Args:
            prs: Presentation avec slides a referencer.
        """
        self._reference_slides = {}
        ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        for idx, slide in enumerate(prs.slides):
            rels_data = []
            for rel in slide.part.rels.values():
                if rel.is_external:
                    continue
                rels_data.append(
                    {
                        "rId": rel.rId,
                        "reltype": rel.reltype,
                        "target_part": rel.target_part,
                    }
                )
            self._reference_slides[idx] = {
                "xml": deepcopy(slide._element),
                "rels": rels_data,
                "layout_name": slide.slide_layout.name,
                "ns_r": ns_r,
            }

    def _clone_reference_slide(self, prs: Presentation, ref_index: int) -> Any:
        """Clone un slide de reference avec tout son design.

        Args:
            prs: Presentation cible.
            ref_index: Index du slide de reference dans le template original.

        Returns:
            Nouveau Slide avec le design clone, ou None si ref_index invalide.
        """
        ref = self._reference_slides.get(ref_index)
        if not ref:
            return None

        # Trouver le layout correspondant
        layout_name = ref["layout_name"]
        layout = self._layout_index.get(layout_name)
        if not layout:
            layout = prs.slide_layouts[_BLANK_LAYOUT_INDEX]

        # Creer un slide vierge
        new_slide = prs.slides.add_slide(layout)

        # Copier shapes, background et color map depuis la reference
        self._copy_shapes_from_ref(ref["xml"], new_slide)
        self._copy_background_from_ref(ref["xml"], new_slide)
        self._copy_color_map_from_ref(ref["xml"], new_slide)

        # Re-creer les relations (images, etc.) et mapper les rIds
        rid_map: dict[str, str] = {}
        for rel_data in ref["rels"]:
            old_rid = rel_data["rId"]
            new_rid = new_slide.part.relate_to(rel_data["target_part"], rel_data["reltype"])
            if old_rid != new_rid:
                rid_map[old_rid] = new_rid

        # Mettre a jour les references rId dans le XML clone
        if rid_map:
            ns_r = ref["ns_r"]
            for attr_suffix in ("embed", "link", "id"):
                full_attr = f"{ns_r}{attr_suffix}"
                for el in new_slide._element.iter():
                    old_val = el.get(full_attr)
                    if old_val and old_val in rid_map:
                        el.set(full_attr, rid_map[old_val])

        return new_slide

    def _copy_shapes_from_ref(self, ref_xml: Any, slide: Any) -> None:
        """Copie les shapes (spTree) de la reference vers le slide."""
        ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        src_csld = ref_xml.find(f"{ns}cSld")
        src_sptree = src_csld.find(f"{ns}spTree") if src_csld is not None else None

        dst_sptree = slide.shapes._spTree
        for child in list(dst_sptree):
            dst_sptree.remove(child)
        if src_sptree is not None:
            for child in src_sptree:
                dst_sptree.append(deepcopy(child))

    def _copy_background_from_ref(self, ref_xml: Any, slide: Any) -> None:
        """Copie le background de la reference vers le slide."""
        ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        src_csld = ref_xml.find(f"{ns}cSld")
        if src_csld is None:
            return
        src_bg = src_csld.find(f"{ns}bg")
        if src_bg is None:
            return
        dst_csld = slide._element.find(f"{ns}cSld")
        existing_bg = dst_csld.find(f"{ns}bg")
        if existing_bg is not None:
            dst_csld.remove(existing_bg)
        dst_csld.insert(0, deepcopy(src_bg))

    def _copy_color_map_from_ref(self, ref_xml: Any, slide: Any) -> None:
        """Copie le clrMapOvr de la reference vers le slide."""
        ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        src_clr = ref_xml.find(f"{ns}clrMapOvr")
        if src_clr is None:
            return
        existing_clr = slide._element.find(f"{ns}clrMapOvr")
        if existing_clr is not None:
            slide._element.remove(existing_clr)
        slide._element.append(deepcopy(src_clr))

    def _replace_cloned_text(
        self, slide: Any, title: str | None, body_texts: list[str] | None = None
    ) -> None:
        """Remplace le texte dans un slide clone.

        Identifie les shapes de texte par taille (area) :
        - Le plus grand = titre
        - Les suivants = contenu (body)

        Args:
            slide: Slide clone.
            title: Nouveau titre (remplace le plus grand text shape).
            body_texts: Textes de remplacement pour les shapes suivantes.
        """
        from pptx.util import Emu  # noqa: PLC0415

        # Collecter les text shapes avec leur taille
        text_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text_frame.text.strip():
                continue
            area = Emu(shape.width).cm * Emu(shape.height).cm
            text_shapes.append((area, shape))

        # Trier par area decroissante (le plus grand = titre)
        text_shapes.sort(key=lambda x: x[0], reverse=True)

        if not text_shapes:
            return

        # Remplacer le titre (plus grand shape)
        if title is not None and text_shapes:
            self._replace_shape_text(text_shapes[0][1], title)

        # Remplacer le body (shapes suivants) et vider les shapes excessifs
        body_shapes = text_shapes[1:]
        if not body_texts:
            # Pas de body demande: vider tous les shapes body
            for _, shape in body_shapes:
                shape.text_frame.clear()
            return
        for i, (_, shape) in enumerate(body_shapes):
            if i >= len(body_texts):
                shape.text_frame.clear()
                continue
            self._replace_shape_text(shape, body_texts[i])

    def _replace_shape_text(self, shape: Any, text: str) -> None:
        """Remplace le texte d'un shape en preservant le style du premier run.

        Args:
            shape: Shape avec text_frame.
            text: Nouveau texte.
        """
        tf = shape.text_frame
        # Sauvegarder le style du premier run
        first_para = tf.paragraphs[0]
        font_props: dict[str, Any] = {}
        if first_para.runs:
            run = first_para.runs[0]
            font_props["bold"] = run.font.bold
            font_props["size"] = run.font.size
            font_props["color"] = run.font.color.rgb if run.font.color.type else None
            font_props["name"] = run.font.name

        # Effacer et remplacer
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        if p.runs and font_props:
            run = p.runs[0]
            if font_props.get("bold") is not None:
                run.font.bold = font_props["bold"]
            if font_props.get("size"):
                run.font.size = font_props["size"]
            if font_props.get("color"):
                run.font.color.rgb = font_props["color"]
            if font_props.get("name"):
                run.font.name = font_props["name"]

    def _get_clone_ref_index(self) -> int | None:
        """Retourne l'index du slide template a cloner pour le slide courant.

        Returns:
            Index du template slide ou None si pas de mapping.
        """
        if not self._template_slide_map:
            return None
        if self._current_slide_idx >= len(self._template_slide_map):
            return None
        return self._template_slide_map[self._current_slide_idx]

    def _process_slides(self, prs: Presentation, slides: list[Any]) -> tuple[int, int]:
        """Cree tous les slides de la presentation.

        Args:
            prs: Presentation python-pptx cible.
            slides: Liste de SlideDef (union discriminee par layout).

        Returns:
            Tuple (total_images, total_charts).
        """
        total_images = 0
        total_charts = 0

        # Layout dispatch table
        slide_handlers: dict[str, Any] = {
            "title_slide": self._create_title_slide,
            "content_bullets": self._create_content_bullets_slide,
            "content_with_image": self._create_content_with_image_slide,
            "section_header": self._create_section_header_slide,
            "two_columns": self._create_two_columns_slide,
            "image_full": self._create_image_full_slide,
            "chart_slide": self._create_chart_slide,
            "closing": self._create_closing_slide,
        }

        for slide_def in slides:
            handler = slide_handlers.get(slide_def.layout)
            if not handler:
                msg = f"Layout inconnu: {slide_def.layout}"
                raise PresentationGenerationError(msg)

            images, charts = handler(prs, slide_def)
            total_images += images
            total_charts += charts
            self._current_slide_idx += 1

        return total_images, total_charts

    # === Layout Handlers ===

    def _create_title_slide(self, prs: Presentation, slide_def: TitleSlideDef) -> tuple[int, int]:
        """Cree un slide de titre."""
        slide = self._add_slide(prs, "title_slide")
        if self._last_slide_cloned:
            self._replace_cloned_text(
                slide, slide_def.title, [slide_def.subtitle] if slide_def.subtitle else None
            )
            self._set_notes(slide, slide_def.notes)
            return 0, 0
        if not self._fill_title_placeholder(slide, slide_def.title, slide_def.title_style):
            self._add_title_textbox(slide, slide_def.title, slide_def.title_style, font_size=44)
        if slide_def.subtitle and not self._fill_subtitle_placeholder(slide, slide_def.subtitle):
            self._add_subtitle_textbox(slide, slide_def.subtitle)
        self._set_notes(slide, slide_def.notes)
        return 0, 0

    def _create_content_bullets_slide(
        self, prs: Presentation, slide_def: ContentBulletsSlideDef
    ) -> tuple[int, int]:
        """Cree un slide avec liste a puces."""
        slide = self._add_slide(prs, "content_bullets")
        if self._last_slide_cloned:
            body = [b.text for b in slide_def.bullets]
            self._replace_cloned_text(slide, slide_def.title, body)
            self._set_notes(slide, slide_def.notes)
            return 0, 0
        if not self._fill_title_placeholder(slide, slide_def.title, slide_def.title_style):
            self._add_title_textbox(slide, slide_def.title, slide_def.title_style)
        if not self._fill_body_placeholder(slide, slide_def.bullets):
            self._add_bullets_textbox(
                slide, slide_def.bullets, left_cm=1.5, top_cm=2.5, width_cm=21.0
            )
        self._set_notes(slide, slide_def.notes)
        return 0, 0

    def _create_content_with_image_slide(
        self, prs: Presentation, slide_def: ContentWithImageSlideDef
    ) -> tuple[int, int]:
        """Cree un slide avec texte et image."""
        slide = self._add_slide(prs, "content_with_image")
        if self._last_slide_cloned:
            body = [b.text for b in slide_def.bullets]
            self._replace_cloned_text(slide, slide_def.title, body)
            self._add_image_to_slide(slide, slide_def.image, left_cm=13.0, top_cm=2.5)
            self._set_notes(slide, slide_def.notes)
            return 1, 0
        if not self._fill_title_placeholder(slide, slide_def.title, slide_def.title_style):
            self._add_title_textbox(slide, slide_def.title, slide_def.title_style)
        if not self._fill_body_placeholder(slide, slide_def.bullets):
            self._add_bullets_textbox(
                slide, slide_def.bullets, left_cm=1.0, top_cm=2.5, width_cm=11.0
            )
        self._add_image_to_slide(slide, slide_def.image, left_cm=13.0, top_cm=2.5)
        self._set_notes(slide, slide_def.notes)
        return 1, 0

    def _create_section_header_slide(
        self, prs: Presentation, slide_def: SectionHeaderSlideDef
    ) -> tuple[int, int]:
        """Cree un slide de transition entre sections."""
        slide = self._add_slide(prs, "section_header")
        if self._last_slide_cloned:
            self._replace_cloned_text(
                slide, slide_def.title, [slide_def.subtitle] if slide_def.subtitle else None
            )
            self._set_notes(slide, slide_def.notes)
            return 0, 0
        if not self._fill_title_placeholder(slide, slide_def.title, slide_def.title_style):
            self._add_title_textbox(
                slide, slide_def.title, slide_def.title_style, top_cm=6.0, font_size=40
            )
        if slide_def.subtitle and not self._fill_subtitle_placeholder(slide, slide_def.subtitle):
            self._add_subtitle_textbox(slide, slide_def.subtitle, top_cm=9.0)
        self._set_notes(slide, slide_def.notes)
        return 0, 0

    def _create_two_columns_slide(
        self, prs: Presentation, slide_def: TwoColumnsSlideDef
    ) -> tuple[int, int]:
        """Cree un slide a deux colonnes."""
        slide = self._add_slide(prs, "two_columns")
        if self._last_slide_cloned:
            body = [b.text for b in slide_def.left_bullets] + [
                b.text for b in slide_def.right_bullets
            ]
            self._replace_cloned_text(slide, slide_def.title, body)
            self._set_notes(slide, slide_def.notes)
            return 0, 0
        if not self._fill_title_placeholder(slide, slide_def.title, slide_def.title_style):
            self._add_title_textbox(slide, slide_def.title, slide_def.title_style)
        left_filled = self._fill_body_placeholder(slide, slide_def.left_bullets, _PH_BODY)
        right_filled = self._fill_body_placeholder(slide, slide_def.right_bullets, _PH_BODY_RIGHT)
        if not left_filled:
            self._add_bullets_textbox(
                slide, slide_def.left_bullets, left_cm=1.0, top_cm=2.5, width_cm=10.5
            )
        if not right_filled:
            self._add_bullets_textbox(
                slide, slide_def.right_bullets, left_cm=12.5, top_cm=2.5, width_cm=10.5
            )
        self._set_notes(slide, slide_def.notes)
        return 0, 0

    def _create_image_full_slide(
        self, prs: Presentation, slide_def: ImageFullSlideDef
    ) -> tuple[int, int]:
        """Cree un slide image pleine page."""
        slide = self._add_slide(prs, "image_full")
        if self._last_slide_cloned:
            self._replace_cloned_text(slide, slide_def.title)
            self._set_notes(slide, slide_def.notes)
            return 0, 0
        self._add_image_to_slide(slide, slide_def.image, left_cm=2.0, top_cm=1.5, full_slide=True)
        if slide_def.title and not self._fill_title_placeholder(
            slide, slide_def.title, slide_def.title_style
        ):
            self._add_title_textbox(slide, slide_def.title, slide_def.title_style, font_size=28)
        if slide_def.caption:
            self._add_caption_textbox(slide, slide_def.caption)
        self._set_notes(slide, slide_def.notes)
        return 1, 0

    def _create_chart_slide(self, prs: Presentation, slide_def: ChartSlideDef) -> tuple[int, int]:
        """Cree un slide avec graphique."""
        slide = self._add_slide(prs, "chart_slide")
        if self._last_slide_cloned:
            self._replace_cloned_text(slide, slide_def.title)
        elif not self._fill_title_placeholder(slide, slide_def.title, slide_def.title_style):
            self._add_title_textbox(slide, slide_def.title, slide_def.title_style)
        try:
            self._add_chart_to_slide(slide, slide_def.chart)
            charts = 1
        except PresentationChartError as exc:
            logger.warning(
                "Graphique ignore '%s': %s",
                slide_def.chart.title or "sans titre",
                str(exc),
            )
            charts = 0
        self._set_notes(slide, slide_def.notes)
        return 0, charts

    def _create_closing_slide(
        self, prs: Presentation, slide_def: ClosingSlideDef
    ) -> tuple[int, int]:
        """Cree un slide de cloture."""
        slide = self._add_slide(prs, "closing")
        if self._last_slide_cloned:
            body = []
            if slide_def.subtitle:
                body.append(slide_def.subtitle)
            if slide_def.bullets:
                body.append("\n".join(b.text for b in slide_def.bullets))
            self._replace_cloned_text(slide, slide_def.title, body if body else None)
            self._set_notes(slide, slide_def.notes)
            return 0, 0
        if not self._fill_title_placeholder(slide, slide_def.title, slide_def.title_style):
            self._add_title_textbox(
                slide, slide_def.title, slide_def.title_style, top_cm=5.0, font_size=40
            )
        if slide_def.subtitle and not self._fill_subtitle_placeholder(slide, slide_def.subtitle):
            self._add_subtitle_textbox(slide, slide_def.subtitle, top_cm=8.0)
        if slide_def.bullets and not self._fill_body_placeholder(slide, slide_def.bullets):
            self._add_bullets_textbox(
                slide, slide_def.bullets, left_cm=6.0, top_cm=10.0, width_cm=12.0
            )
        self._set_notes(slide, slide_def.notes)
        return 0, 0

    # === Helpers ===

    def _add_blank_slide(self, prs: Presentation) -> Slide:
        """Ajoute un slide vierge a la presentation."""
        layout = prs.slide_layouts[_BLANK_LAYOUT_INDEX]
        return prs.slides.add_slide(layout)

    def _build_layout_index(self, prs: Presentation) -> dict[str, Any]:
        """Construit un index nom -> SlideLayout pour la presentation.

        Args:
            prs: Instance Presentation.

        Returns:
            Dict {nom_layout: objet SlideLayout}.
        """
        return {layout.name: layout for layout in prs.slide_layouts}

    def _find_layout(self, layout_type: str) -> Any:
        """Trouve le meilleur layout natif pour un type de slide.

        Args:
            layout_type: Type de notre layout (ex: "title_slide").

        Returns:
            SlideLayout natif ou None si aucun match.
        """
        candidates = _LAYOUT_NAME_CANDIDATES.get(layout_type, [])
        for name in candidates:
            layout = self._layout_index.get(name)
            if layout is not None:
                return layout
        return None

    def _add_slide(self, prs: Presentation, layout_type: str) -> Slide:
        """Ajoute un slide en utilisant clonage, layout natif ou Blank.

        Priorite:
        1. Clone d'un slide de reference (si template_slide_map defini)
        2. Layout natif du template (si template avec layouts nommes)
        3. Blank + textboxes manuels (fallback)

        Args:
            prs: Instance Presentation.
            layout_type: Type de notre layout.

        Returns:
            Slide ajoute.
        """
        self._last_slide_cloned = False

        # Priorite 1: clonage d'un slide de reference
        ref_index = self._get_clone_ref_index()
        if ref_index is not None:
            cloned = self._clone_reference_slide(prs, ref_index)
            if cloned:
                self._last_slide_cloned = True
                return cloned

        if not self._has_template:
            return self._add_blank_slide(prs)

        # Priorite 2: layout natif
        native_layout = self._find_layout(layout_type)
        if not native_layout:
            return self._add_blank_slide(prs)

        return prs.slides.add_slide(native_layout)

    def _get_placeholder(self, slide: Slide, idx: int) -> Any | None:
        """Recupere un placeholder par son index.

        Args:
            slide: Slide cible.
            idx: Index du placeholder.

        Returns:
            Placeholder ou None si non trouve.
        """
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _fill_title_placeholder(
        self, slide: Slide, title: str, style: TextStyle | None = None
    ) -> bool:
        """Remplit le placeholder titre natif s'il existe.

        Args:
            slide: Slide cible.
            title: Texte du titre.
            style: Style optionnel.

        Returns:
            True si le placeholder a ete rempli, False sinon.
        """
        ph = self._get_placeholder(slide, _PH_TITLE)
        if not ph:
            return False
        if not ph.has_text_frame:
            return False
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        if p.runs and style:
            run = p.runs[0]
            if style.bold:
                run.font.bold = True
            if style.italic:
                run.font.italic = True
            if style.font_size:
                run.font.size = Pt(style.font_size)
            if style.font_color:
                run.font.color.rgb = RGBColor.from_string(style.font_color)
        return True

    def _fill_body_placeholder(
        self, slide: Slide, bullets: list[BulletItem], idx: int = _PH_BODY
    ) -> bool:
        """Remplit un placeholder body avec des puces.

        Args:
            slide: Slide cible.
            bullets: Liste de puces.
            idx: Index du placeholder (default: 1).

        Returns:
            True si le placeholder a ete rempli, False sinon.
        """
        ph = self._get_placeholder(slide, idx)
        if not ph:
            return False
        if not ph.has_text_frame:
            return False
        tf = ph.text_frame
        tf.clear()
        self._populate_bullets(tf, bullets)
        return True

    def _fill_subtitle_placeholder(self, slide: Slide, subtitle: str) -> bool:
        """Remplit le placeholder sous-titre natif s'il existe.

        Args:
            slide: Slide cible.
            subtitle: Texte du sous-titre.

        Returns:
            True si le placeholder a ete rempli, False sinon.
        """
        ph = self._get_placeholder(slide, _PH_BODY)
        if not ph:
            return False
        if not ph.has_text_frame:
            return False
        tf = ph.text_frame
        tf.clear()
        tf.paragraphs[0].text = subtitle
        return True

    def _add_title_textbox(
        self,
        slide: Slide,
        title: str,
        style: TextStyle | None = None,
        left_cm: float = 1.5,
        top_cm: float = 0.5,
        width_cm: float = 21.0,
        height_cm: float = 2.0,
        font_size: int = 32,
    ) -> None:
        """Ajoute un titre en textbox sur le slide."""
        txbox = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT

        effective_size = font_size
        if style and style.font_size:
            effective_size = style.font_size

        run = p.runs[0]
        run.font.size = Pt(effective_size)
        run.font.bold = style.bold if style else False
        if style and style.italic:
            run.font.italic = True
        if style and style.font_color:
            run.font.color.rgb = RGBColor.from_string(style.font_color)

    def _add_subtitle_textbox(
        self,
        slide: Slide,
        subtitle: str,
        top_cm: float = 3.5,
    ) -> None:
        """Ajoute un sous-titre sur le slide."""
        txbox = slide.shapes.add_textbox(Cm(1.5), Cm(top_cm), Cm(21.0), Cm(1.5))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    def _add_caption_textbox(self, slide: Slide, caption: str) -> None:
        """Ajoute un texte de legende en bas du slide."""
        txbox = slide.shapes.add_textbox(Cm(1.5), Cm(16.5), Cm(21.0), Cm(1.0))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(12)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    def _add_bullets_textbox(
        self,
        slide: Slide,
        bullets: list[BulletItem],
        left_cm: float,
        top_cm: float,
        width_cm: float,
        height_cm: float = 12.0,
    ) -> None:
        """Ajoute une zone de texte avec puces."""
        txbox = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
        tf = txbox.text_frame
        tf.word_wrap = True
        self._populate_bullets(tf, bullets)

    def _populate_bullets(self, tf: TextFrame, bullets: list[BulletItem]) -> None:
        """Remplit un TextFrame avec des puces."""
        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()

            p.text = bullet.text
            p.level = bullet.level
            p.space_after = Pt(6)

            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(18 - bullet.level * 2)
                run.font.bold = bullet.bold

    def _add_image_to_slide(
        self,
        slide: Slide,
        image_def: ImageDef,
        left_cm: float = 2.0,
        top_cm: float = 2.0,
        full_slide: bool = False,
    ) -> None:
        """Ajoute une image au slide.

        Args:
            slide: Slide cible.
            image_def: Definition de l'image.
            left_cm: Position gauche en cm.
            top_cm: Position haute en cm.
            full_slide: Si True, utilise des dimensions plus grandes.

        Raises:
            PresentationImageNotFoundError: Si l'image est introuvable.
        """
        image_path = self._resolve_image_path(image_def.filename)

        width = Cm(image_def.width_cm) if image_def.width_cm else None
        height = Cm(image_def.height_cm) if image_def.height_cm else None

        if not width and not height:
            width = Cm(20.0) if full_slide else Cm(10.0)

        slide.shapes.add_picture(str(image_path), Cm(left_cm), Cm(top_cm), width, height)

    def _resolve_image_path(self, filename: str) -> Path:
        """Resout le chemin d'une image avec securite.

        Args:
            filename: Nom du fichier image.

        Returns:
            Chemin absolu vers l'image.

        Raises:
            PresentationImageNotFoundError: Si l'image n'existe pas.
            PresentationGenerationError: Si path traversal detecte.
        """
        if ".." in filename or "/" in filename or "\\" in filename:
            raise PresentationGenerationError(
                message=f"Nom d'image invalide: {filename}",
                code="INVALID_IMAGE_NAME",
                suggestion="Le nom d'image ne doit pas contenir '..' , '/' ou '\\'.",
                parameter="image",
                retry=True,
            )

        image_path = self._images_path / filename
        resolved = image_path.resolve()
        if not resolved.is_relative_to(self._images_path.resolve()):
            raise PresentationGenerationError(
                message=f"Path traversal detecte: {filename}",
                code="PATH_TRAVERSAL",
                suggestion="Utiliser un nom de fichier simple sans chemin.",
                parameter="image",
                retry=True,
            )

        if not image_path.exists():
            available = self._list_available_images()
            raise PresentationImageNotFoundError(filename, available)

        return image_path

    def _list_available_images(self) -> list[str]:
        """Liste les images disponibles."""
        if not self._images_path.exists():
            return []
        extensions = ("*.png", "*.jpg", "*.jpeg", "*.gif", "*.bmp", "*.svg")
        images: list[str] = []
        for ext in extensions:
            images.extend(f.name for f in self._images_path.glob(ext))
        return sorted(images)

    def _add_chart_to_slide(self, slide: Slide, chart_def: PresentationChartDef) -> None:
        """Ajoute un graphique au slide.

        Raises:
            PresentationChartError: En cas d'echec.
        """
        try:
            chart_type = _CHART_TYPE_MAP.get(chart_def.chart_type)
            if chart_type is None:
                msg = f"Type de graphique non supporte: {chart_def.chart_type}"
                raise PresentationGenerationError(msg)

            chart_data = CategoryChartData()
            chart_data.categories = chart_def.categories
            for series in chart_def.series:
                chart_data.add_series(series.name, series.values)

            chart_frame = slide.shapes.add_chart(
                chart_type,
                Cm(2.0),
                Cm(3.0),
                Cm(20.0),
                Cm(13.0),
                chart_data,
            )

            if chart_def.title:
                chart_frame.chart.has_title = True
                chart_frame.chart.chart_title.text_frame.text = chart_def.title

        except (PresentationChartError, PresentationGenerationError):
            raise
        except Exception as exc:
            raise PresentationChartError(
                chart_title=chart_def.title,
                reason=str(exc),
            ) from exc

    def _set_notes(self, slide: Slide, notes: str | None) -> None:
        """Definit les notes du presentateur sur un slide."""
        if not notes:
            return
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    # === File Management ===

    def save_and_verify(self, prs: Presentation, file_path: Path, filename: str) -> int:
        """Sauvegarde la presentation et verifie la taille.

        Returns:
            Taille du fichier en octets.

        Raises:
            PresentationOutputTooLargeError: Si le fichier depasse la taille max.
        """
        prs.save(str(file_path))
        file_size = file_path.stat().st_size
        size_mb = file_size / (1024 * 1024)
        if size_mb > self._max_output_size_mb:
            file_path.unlink()
            raise PresentationOutputTooLargeError(
                filename=filename,
                size_mb=size_mb,
                max_size_mb=self._max_output_size_mb,
            )
        return file_size

    def sanitize_filename(self, filename: str) -> str:
        """Securise le nom de fichier (path traversal prevention).

        Args:
            filename: Nom de fichier a securiser.

        Returns:
            Nom de fichier securise.

        Raises:
            PresentationGenerationError: Si le nom est invalide ou dangereux.
        """
        if ".." in filename or "/" in filename or "\\" in filename:
            raise PresentationGenerationError(
                message=f"Nom de fichier invalide: {filename}",
                code="INVALID_FILENAME",
                suggestion="Le nom de fichier ne doit pas contenir '..' , '/' ou '\\'.",
                parameter="filename",
                retry=True,
            )

        resolved = (self._output_path / filename).resolve()
        if not resolved.is_relative_to(self._output_path.resolve()):
            raise PresentationGenerationError(
                message=f"Path traversal detecte: {filename}",
                code="PATH_TRAVERSAL",
                suggestion="Utiliser un nom de fichier simple sans chemin.",
                parameter="filename",
                retry=True,
            )

        return filename

    def ensure_output_dir(self) -> None:
        """Cree le repertoire OUTPUT_PATH s'il n'existe pas."""
        self._output_path.mkdir(parents=True, exist_ok=True)
