# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Assistance Micro Design
"""Editeur de presentations PowerPoint (.pptx) existantes."""

from __future__ import annotations

import asyncio
import logging
from collections.abc import Callable
from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt


if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide

from src.core.config import settings
from src.core.exceptions import (
    PresentationChartError,
    PresentationFileNotFoundError,
    PresentationGenerationError,
    PresentationSlideNotFoundError,
)
from src.models.api import EditPresentationParams, EditPresentationResult
from src.models.presentation_edit import (
    AddImageOp,
    AddSlideOp,
    DeleteSlideOp,
    ReorderSlideOp,
    ReplaceImageOp,
    SetBackgroundOp,
    UpdateBulletsOp,
    UpdateChartOp,
    UpdateNotesOp,
    UpdateSubtitleOp,
    UpdateTitleOp,
)
from src.services.presentation.generator import PresentationGenerator
from src.services.presentation.shape_finders import (
    find_bullets_shape,
    find_subtitle_shape,
    find_title_shape,
)


logger = logging.getLogger(__name__)


class PresentationEditor:
    """Editeur de presentations PowerPoint existantes.

    Utilise PresentationGenerator en composition pour reutiliser
    les helpers (images, charts, bullets).

    Attributes:
        _output_path: Repertoire contenant les fichiers a editer.
        _generator: Instance PresentationGenerator pour reutiliser les helpers.
    """

    def __init__(self, output_path: Path | None = None) -> None:
        """Initialise l'editeur.

        Args:
            output_path: Repertoire des fichiers (defaut: settings.OUTPUT_PATH).
        """
        self._output_path = Path(output_path or settings.OUTPUT_PATH)
        self._generator = PresentationGenerator(output_path=self._output_path)

    async def edit(self, params: EditPresentationParams) -> EditPresentationResult:
        """Point d'entree principal. Edite le fichier pptx.

        Args:
            params: Parametres d'edition (filename + operations).

        Returns:
            Resultat avec stats des operations.

        Raises:
            PresentationFileNotFoundError: Si le fichier n'existe pas.
            PresentationSlideNotFoundError: Si un slide reference n'existe pas.
        """
        return await asyncio.to_thread(self._edit_sync, params)

    def _edit_sync(self, params: EditPresentationParams) -> EditPresentationResult:
        """Edition synchrone du fichier (appelee via to_thread).

        Returns:
            Resultat avec stats des operations appliquees.
        """
        safe_filename = self._generator.sanitize_filename(params.filename)
        file_path = self._output_path / safe_filename
        if not file_path.exists():
            raise PresentationFileNotFoundError(safe_filename)

        prs = Presentation(str(file_path))
        applied, skipped = self._apply_operations(prs, params.operations)
        file_size = self._generator.save_and_verify(prs, file_path, safe_filename)

        logger.info(
            "Presentation editee: %s (%d ops, %d ignorees, %d octets)",
            safe_filename,
            applied,
            skipped,
            file_size,
        )

        return EditPresentationResult(
            file_path=str(file_path),
            filename=safe_filename,
            operations_applied=applied,
            operations_skipped=skipped,
            file_size_bytes=file_size,
        )

    def _apply_operations(self, prs: Presentation, operations: list[Any]) -> tuple[int, int]:
        """Applique les operations en sequence via dispatch dict.

        Les PresentationChartError sont catchees et comptees comme skipped.

        Returns:
            Tuple (applied, skipped).
        """
        op_handlers: dict[str, Callable[..., None]] = {
            "update_title": self._op_update_title,
            "update_subtitle": self._op_update_subtitle,
            "update_bullets": self._op_update_bullets,
            "add_slide": self._op_add_slide,
            "delete_slide": self._op_delete_slide,
            "reorder_slide": self._op_reorder_slide,
            "replace_image": self._op_replace_image,
            "add_image": self._op_add_image,
            "update_notes": self._op_update_notes,
            "update_chart": self._op_update_chart,
            "set_background": self._op_set_background,
        }
        applied = 0
        skipped = 0
        for op in operations:
            handler = op_handlers.get(op.op)
            if not handler:
                msg = f"Operation inconnue: {op.op}"
                raise PresentationGenerationError(msg)
            try:
                handler(prs, op)
                applied += 1
            except PresentationChartError as exc:
                logger.warning("Operation chart ignoree: %s", exc)
                skipped += 1
        return applied, skipped

    # === Helpers ===

    def _get_slide(self, prs: Presentation, index: int) -> Slide:
        """Recupere un slide par index ou leve PresentationSlideNotFoundError."""
        total = len(prs.slides)
        if index < 0 or index >= total:
            raise PresentationSlideNotFoundError(index, total)
        return prs.slides[index]

    def _find_title_shape(self, slide: Slide) -> BaseShape | None:
        """Trouve le premier shape contenant du texte (titre presume)."""
        return find_title_shape(slide)

    def _find_subtitle_shape(self, slide: Slide) -> BaseShape | None:
        """Trouve le deuxieme shape contenant du texte (sous-titre presume)."""
        return find_subtitle_shape(slide)

    def _find_bullets_shape(self, slide: Slide) -> BaseShape | None:
        """Trouve le shape contenant des puces (le plus grand textbox apres le titre)."""
        return find_bullets_shape(slide)

    # === Operations ===

    def _op_update_title(self, prs: Presentation, op: UpdateTitleOp) -> None:
        """Modifier le titre d'un slide."""
        slide = self._get_slide(prs, op.slide_index)
        shape = self._find_title_shape(slide)
        if not shape:
            # Aucun shape texte existant — creer un textbox titre a la place
            self._generator._add_title_textbox(slide, op.title, op.style)
            return
        tf = shape.text_frame
        tf.paragraphs[0].text = op.title
        if op.style and tf.paragraphs[0].runs:
            run = tf.paragraphs[0].runs[0]
            run.font.bold = op.style.bold
            if op.style.font_size:
                run.font.size = Pt(op.style.font_size)
            if op.style.italic:
                run.font.italic = True
            if op.style.font_color:
                run.font.color.rgb = RGBColor.from_string(op.style.font_color)

    def _op_update_subtitle(self, prs: Presentation, op: UpdateSubtitleOp) -> None:
        """Modifier le sous-titre d'un slide."""
        slide = self._get_slide(prs, op.slide_index)
        shape = self._find_subtitle_shape(slide)
        if not shape:
            self._generator._add_subtitle_textbox(slide, op.subtitle)
            return
        shape.text_frame.paragraphs[0].text = op.subtitle

    def _op_update_bullets(self, prs: Presentation, op: UpdateBulletsOp) -> None:
        """Remplacer les puces d'un slide."""
        slide = self._get_slide(prs, op.slide_index)
        shape = self._find_bullets_shape(slide)
        if not shape:
            self._generator._add_bullets_textbox(
                slide, op.bullets, left_cm=1.5, top_cm=2.5, width_cm=21.0
            )
            return
        # Vider pour eviter de dupliquer les puces existantes
        tf = shape.text_frame
        tf.clear()
        self._generator._populate_bullets(tf, op.bullets)

    def _op_add_slide(self, prs: Presentation, op: AddSlideOp) -> None:
        """Ajouter un slide."""
        # Deleguer au generator pour reutiliser la creation par layout
        self._generator._process_slides(prs, [op.slide])

        # Deplacer le slide ajoute en fin vers la position cible
        if op.at_index is not None:
            total = len(prs.slides)
            if op.at_index < total - 1:
                self._move_slide(prs, total - 1, op.at_index)

    def _op_delete_slide(self, prs: Presentation, op: DeleteSlideOp) -> None:
        """Supprimer un slide."""
        self._get_slide(prs, op.slide_index)  # Validates index
        sld_id_lst = self._get_sld_id_lst(prs)
        sld_ids = list(sld_id_lst)
        sld_id_lst.remove(sld_ids[op.slide_index])

    def _op_reorder_slide(self, prs: Presentation, op: ReorderSlideOp) -> None:
        """Deplacer un slide."""
        total = len(prs.slides)
        if op.from_index >= total:
            raise PresentationSlideNotFoundError(op.from_index, total)
        if op.to_index >= total:
            raise PresentationSlideNotFoundError(op.to_index, total)
        self._move_slide(prs, op.from_index, op.to_index)

    def _get_sld_id_lst(self, prs: Presentation) -> Any:
        """Recupere l'element XML sldIdLst (API privee: pas de methode publique pour manipuler l'ordre)."""
        return prs.part._element.find(qn("p:sldIdLst"))

    def _move_slide(self, prs: Presentation, from_idx: int, to_idx: int) -> None:
        """Deplace un slide d'une position a une autre."""
        sld_id_lst = self._get_sld_id_lst(prs)
        sld_ids = list(sld_id_lst)
        element = sld_ids[from_idx]
        sld_id_lst.remove(element)
        sld_ids_updated = list(sld_id_lst)
        if to_idx >= len(sld_ids_updated):
            sld_id_lst.append(element)
            return
        sld_ids_updated[to_idx].addprevious(element)

    def _op_replace_image(self, prs: Presentation, op: ReplaceImageOp) -> None:
        """Remplacer l'image d'un slide."""
        slide = self._get_slide(prs, op.slide_index)
        # python-pptx n'expose pas de methode publique pour supprimer un shape
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sp = shape._element
                sp.getparent().remove(sp)
                break
        self._generator._add_image_to_slide(slide, op.image)

    def _op_add_image(self, prs: Presentation, op: AddImageOp) -> None:
        """Ajouter une image a un slide."""
        slide = self._get_slide(prs, op.slide_index)
        self._generator._add_image_to_slide(slide, op.image)

    def _op_update_notes(self, prs: Presentation, op: UpdateNotesOp) -> None:
        """Modifier les notes du presentateur."""
        slide = self._get_slide(prs, op.slide_index)
        self._generator._set_notes(slide, op.notes)

    def _op_update_chart(self, prs: Presentation, op: UpdateChartOp) -> None:
        """Remplacer le graphique d'un slide."""
        slide = self._get_slide(prs, op.slide_index)
        # python-pptx n'expose pas de methode publique pour supprimer un shape
        for shape in list(slide.shapes):
            if shape.has_chart:
                sp = shape._element
                sp.getparent().remove(sp)
                break
        self._generator._add_chart_to_slide(slide, op.chart)

    def _op_set_background(self, prs: Presentation, op: SetBackgroundOp) -> None:
        """Definir la couleur de fond d'un slide."""
        slide = self._get_slide(prs, op.slide_index)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(op.color)
