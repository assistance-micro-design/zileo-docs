# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Assistance Micro Design
"""Service d'inspection de templates PowerPoint (.pptx)."""

from __future__ import annotations

import asyncio
import contextlib
import logging
from pathlib import Path
from typing import Any

from src.core.config import settings


logger = logging.getLogger(__name__)

_MAX_LAYOUTS = 20


class TemplateInspector:
    """Inspecte la structure d'un template PowerPoint.

    Extrait les layouts, placeholders, theme et dimensions
    pour permettre au LLM de comprendre la structure du template.

    Attributes:
        _templates_path: Repertoire contenant les templates.
    """

    def __init__(self, templates_path: Path | None = None) -> None:
        """Initialise l'inspecteur.

        Args:
            templates_path: Repertoire des templates (defaut: settings.TEMPLATES_PPTX_PATH).
        """
        self._templates_path = Path(templates_path or settings.TEMPLATES_PPTX_PATH)

    async def inspect(self, template: str) -> dict[str, Any]:
        """Inspecte un template et retourne sa structure.

        Args:
            template: Nom du fichier template .pptx.

        Returns:
            Structure du template ou dict avec "error" en cas de probleme.
        """
        if ".." in template or "/" in template or "\\" in template:
            return {"error": f"Nom de template invalide: {template}"}

        template_path = self._templates_path / template
        if not template_path.exists():
            return {
                "error": f"Template introuvable: {template}",
                "available_templates": self._list_templates(),
            }

        return await asyncio.to_thread(self._inspect_sync, template_path, template)

    def _inspect_sync(self, template_path: Path, template_name: str) -> dict[str, Any]:
        """Inspection synchrone du template.

        Args:
            template_path: Chemin absolu du fichier template.
            template_name: Nom du fichier template.

        Returns:
            Structure du template.
        """
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Emu  # noqa: PLC0415

        try:
            prs = Presentation(str(template_path))
        except Exception as exc:
            return {"error": f"Fichier illisible: {exc}", "template": template_name}

        layouts = self._extract_layouts(prs, Emu)
        theme = self._extract_theme(prs)
        total_layouts = len(list(prs.slide_layouts))
        slides = self._extract_slides(prs, Emu)

        return {
            "template": template_name,
            "file_size_bytes": template_path.stat().st_size,
            "slide_width_cm": round(Emu(prs.slide_width).cm, 2),
            "slide_height_cm": round(Emu(prs.slide_height).cm, 2),
            "theme": theme,
            "total_layouts": total_layouts,
            "layouts": layouts,
            "existing_slides_count": len(prs.slides),
            "slides": slides,
            "hint": (
                f"Ce template a {len(slides)} slides et {total_layouts} layouts. "
                "Utilisez template_slide_map pour cloner des slides existants "
                "en preservant tout le design (fonds, images, formes, polices). "
                "Chaque slide[i] clone le template slide template_slide_map[i]. "
                "Les text shapes sont tries par taille: le plus grand = titre, les suivants = body. "
                "Sans template_slide_map, les layouts natifs du template sont utilises."
            ),
        }

    def _extract_layouts(self, prs: Any, emu_cls: Any) -> list[dict[str, Any]]:
        """Extrait les layouts et leurs placeholders.

        Args:
            prs: Instance Presentation python-pptx.
            emu_cls: Classe Emu pour conversion d'unites.

        Returns:
            Liste des layouts avec placeholders.
        """
        layouts: list[dict[str, Any]] = []
        for idx, layout in enumerate(prs.slide_layouts):
            if idx >= _MAX_LAYOUTS:
                break
            placeholders = self._extract_placeholders(layout, emu_cls)
            layouts.append(
                {
                    "index": idx,
                    "name": layout.name,
                    "placeholders": placeholders,
                }
            )
        return layouts

    def _extract_placeholders(self, layout: Any, emu_cls: Any) -> list[dict[str, Any]]:
        """Extrait les placeholders d'un layout.

        Args:
            layout: Objet SlideLayout python-pptx.
            emu_cls: Classe Emu pour conversion d'unites.

        Returns:
            Liste des placeholders avec index, nom, type et dimensions.
        """
        placeholders: list[dict[str, Any]] = []
        for ph in layout.placeholders:
            placeholders.append(
                {
                    "idx": ph.placeholder_format.idx,
                    "name": ph.name,
                    "type": str(ph.placeholder_format.type),
                    "left_cm": round(emu_cls(ph.left).cm, 2),
                    "top_cm": round(emu_cls(ph.top).cm, 2),
                    "width_cm": round(emu_cls(ph.width).cm, 2),
                    "height_cm": round(emu_cls(ph.height).cm, 2),
                }
            )
        return placeholders

    def _extract_theme(self, prs: Any) -> dict[str, Any]:
        """Extrait les informations du theme.

        Args:
            prs: Instance Presentation python-pptx.

        Returns:
            Dictionnaire avec nom, color_scheme, font_major, font_minor.
        """
        theme: dict[str, Any] = {
            "name": None,
            "color_scheme": None,
            "font_major": None,
            "font_minor": None,
        }

        if not prs.slide_masters:
            return theme

        master = prs.slide_masters[0]
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        with contextlib.suppress(AttributeError, TypeError, IndexError):
            theme_el = master.element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
            )
            if theme_el is not None:
                theme["name"] = theme_el.get("name")

        with contextlib.suppress(AttributeError, TypeError, IndexError):
            font_scheme = master.element.find(".//a:fontScheme", ns)
            if font_scheme is not None:
                major = font_scheme.find(".//a:majorFont/a:latin", ns)
                minor = font_scheme.find(".//a:minorFont/a:latin", ns)
                if major is not None:
                    theme["font_major"] = major.get("typeface")
                if minor is not None:
                    theme["font_minor"] = minor.get("typeface")

        with contextlib.suppress(AttributeError, TypeError, IndexError):
            clr_scheme = master.element.find(".//a:clrScheme", ns)
            if clr_scheme is not None:
                theme["color_scheme"] = clr_scheme.get("name")

        return theme

    def _extract_slides(self, prs: Any, emu_cls: Any) -> list[dict[str, Any]]:
        """Extrait le contenu des slides existants du template.

        Args:
            prs: Instance Presentation python-pptx.
            emu_cls: Classe Emu pour conversion d'unites.

        Returns:
            Liste des slides avec leurs text shapes tries par taille.
        """
        slides: list[dict[str, Any]] = []
        for idx, slide in enumerate(prs.slides):
            text_shapes = self._extract_text_shapes(slide, emu_cls)
            empty_text_shapes = self._extract_empty_text_shapes(slide, emu_cls)
            images = self._extract_slide_images(slide, emu_cls)
            groups = self._extract_groups(slide, emu_cls)
            decorations = self._extract_decorations(slide, emu_cls)
            notes = self._extract_slide_notes(slide)
            slides.append(
                {
                    "index": idx,
                    "layout": slide.slide_layout.name,
                    "total_shapes": len(slide.shapes),
                    "text_shapes": text_shapes,
                    "empty_text_shapes": empty_text_shapes,
                    "images": images,
                    "groups": groups,
                    "decorations": decorations,
                    "notes": notes,
                }
            )
        return slides

    def _extract_text_shapes(self, slide: Any, emu_cls: Any) -> list[dict[str, Any]]:
        """Extrait les text shapes d'un slide tries par area decroissante.

        Args:
            slide: Objet Slide python-pptx.
            emu_cls: Classe Emu pour conversion d'unites.

        Returns:
            Liste des shapes avec text, font, position et role estime.
        """
        shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            width_cm = round(emu_cls(shape.width).cm, 1)
            height_cm = round(emu_cls(shape.height).cm, 1)
            area = round(width_cm * height_cm, 1)
            font_info = self._extract_font_info(shape)
            shapes.append(
                {
                    "text": text[:100],
                    "word_count": len(text.split()),
                    "font": font_info,
                    "width_cm": width_cm,
                    "height_cm": height_cm,
                    "area_cm2": area,
                    "left_cm": round(emu_cls(shape.left).cm, 1),
                    "top_cm": round(emu_cls(shape.top).cm, 1),
                }
            )
        shapes.sort(key=lambda s: s["area_cm2"], reverse=True)
        for i, shape in enumerate(shapes):
            shape["role"] = "title" if i == 0 else f"body_{i}"
        return shapes

    def _extract_font_info(self, shape: Any) -> dict[str, Any]:
        """Extrait les infos de police du premier run du shape.

        Args:
            shape: Shape avec text_frame.

        Returns:
            Dict avec size_pt, bold, name, color.
        """
        info: dict[str, Any] = {
            "size_pt": None,
            "bold": None,
            "name": None,
            "color": None,
        }
        for para in shape.text_frame.paragraphs:
            if not para.runs:
                continue
            run = para.runs[0]
            if run.font.size:
                info["size_pt"] = round(run.font.size.pt, 1)
            info["bold"] = run.font.bold
            info["name"] = run.font.name
            with contextlib.suppress(AttributeError, TypeError):
                if run.font.color and run.font.color.type:
                    info["color"] = str(run.font.color.rgb)
            return info
        return info

    def _extract_slide_images(self, slide: Any, emu_cls: Any) -> list[dict[str, Any]]:
        """Extrait les images d'un slide.

        Args:
            slide: Objet Slide python-pptx.
            emu_cls: Classe Emu pour conversion d'unites.

        Returns:
            Liste des images avec dimensions et position.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

        images: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            content_type = None
            with contextlib.suppress(AttributeError):
                content_type = shape.image.content_type
            images.append(
                {
                    "name": shape.name,
                    "content_type": content_type,
                    "width_cm": round(emu_cls(shape.width).cm, 1),
                    "height_cm": round(emu_cls(shape.height).cm, 1),
                    "left_cm": round(emu_cls(shape.left).cm, 1),
                    "top_cm": round(emu_cls(shape.top).cm, 1),
                }
            )
        return images

    def _extract_decorations(self, slide: Any, emu_cls: Any) -> list[dict[str, Any]]:
        """Extrait les formes decoratives (non-texte, non-image) d'un slide.

        Args:
            slide: Objet Slide python-pptx.
            emu_cls: Classe Emu pour conversion d'unites.

        Returns:
            Liste des formes avec type, dimensions et position.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

        decorations: list[dict[str, Any]] = []
        skip_types = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP}
        for shape in slide.shapes:
            if shape.shape_type in skip_types:
                continue
            if shape.has_text_frame and shape.text_frame.text.strip():
                continue
            width = round(emu_cls(shape.width).cm, 1)
            height = round(emu_cls(shape.height).cm, 1)
            if width < 0.5 and height < 0.5:
                continue
            shape_type = str(shape.shape_type) if shape.shape_type else "unknown"
            decorations.append(
                {
                    "name": shape.name,
                    "type": shape_type,
                    "width_cm": width,
                    "height_cm": height,
                    "left_cm": round(emu_cls(shape.left).cm, 1),
                    "top_cm": round(emu_cls(shape.top).cm, 1),
                }
            )
        return decorations

    def _extract_empty_text_shapes(self, slide: Any, emu_cls: Any) -> list[dict[str, Any]]:
        """Extrait les shapes avec text_frame vide (zones de saisie potentielles).

        Args:
            slide: Objet Slide python-pptx.
            emu_cls: Classe Emu pour conversion d'unites.

        Returns:
            Liste des shapes texte vides avec dimensions et position.
        """
        shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.text_frame.text.strip():
                continue
            width = round(emu_cls(shape.width).cm, 1)
            height = round(emu_cls(shape.height).cm, 1)
            if width < 0.5 and height < 0.5:
                continue
            shapes.append(
                {
                    "name": shape.name,
                    "width_cm": width,
                    "height_cm": height,
                    "left_cm": round(emu_cls(shape.left).cm, 1),
                    "top_cm": round(emu_cls(shape.top).cm, 1),
                }
            )
        return shapes

    def _extract_groups(self, slide: Any, emu_cls: Any) -> list[dict[str, Any]]:
        """Extrait les groupes de shapes (icones, compositions).

        Args:
            slide: Objet Slide python-pptx.
            emu_cls: Classe Emu pour conversion d'unites.

        Returns:
            Liste des groupes avec nombre d'enfants et dimensions.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

        groups: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                continue
            n_children = len(shape.shapes) if hasattr(shape, "shapes") else 0
            groups.append(
                {
                    "name": shape.name,
                    "children_count": n_children,
                    "width_cm": round(emu_cls(shape.width).cm, 1),
                    "height_cm": round(emu_cls(shape.height).cm, 1),
                    "left_cm": round(emu_cls(shape.left).cm, 1),
                    "top_cm": round(emu_cls(shape.top).cm, 1),
                }
            )
        return groups

    def _extract_slide_notes(self, slide: Any) -> str | None:
        """Extrait les notes du presentateur d'un slide.

        Args:
            slide: Objet Slide python-pptx.

        Returns:
            Texte des notes ou None.
        """
        if not slide.has_notes_slide:
            return None
        if not slide.notes_slide.notes_text_frame:
            return None
        text: str = slide.notes_slide.notes_text_frame.text
        if not text.strip():
            return None
        return text.strip()

    def _list_templates(self) -> list[str]:
        """Liste les templates disponibles."""
        if not self._templates_path.exists():
            return []
        return sorted(f.name for f in self._templates_path.glob("*.pptx"))
