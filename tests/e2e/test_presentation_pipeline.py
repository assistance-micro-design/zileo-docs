# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Assistance Micro Design
"""Tests E2E pour le pipeline de presentation PowerPoint: create -> edit -> verify."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from src.models.api import (
    CreatePresentationParams,
    EditPresentationParams,
)
from src.models.presentation_edit import (
    AddSlideOp,
    DeleteSlideOp,
    SetBackgroundOp,
    UpdateBulletsOp,
    UpdateNotesOp,
    UpdateTitleOp,
)
from src.models.presentation_generation import (
    BulletItem,
    ChartSeriesDef,
    ChartSlideDef,
    ClosingSlideDef,
    ContentBulletsSlideDef,
    ContentWithImageSlideDef,
    ImageDef,
    ImageFullSlideDef,
    PresentationChartDef,
    SectionHeaderSlideDef,
    TextStyle,
    TitleSlideDef,
    TwoColumnsSlideDef,
)
from src.services.presentation.editor import PresentationEditor
from src.services.presentation.generator import PresentationGenerator


@pytest.fixture
def tmp_output(tmp_path: Path) -> Path:
    output = tmp_path / "output"
    output.mkdir()
    return output


@pytest.fixture
def tmp_images(tmp_path: Path) -> Path:
    from tests.conftest import create_minimal_png

    images = tmp_path / "images"
    images.mkdir()
    create_minimal_png(images / "logo.png")
    create_minimal_png(images / "diagram.png")
    return images


@pytest.fixture
def tmp_templates(tmp_path: Path) -> Path:
    templates = tmp_path / "templates"
    templates.mkdir()
    return templates


@pytest.fixture
def generator(tmp_output: Path, tmp_images: Path, tmp_templates: Path) -> PresentationGenerator:
    return PresentationGenerator(
        output_path=tmp_output,
        images_path=tmp_images,
        templates_path=tmp_templates,
    )


@pytest.fixture
def editor(tmp_output: Path) -> PresentationEditor:
    return PresentationEditor(output_path=tmp_output)


@pytest.mark.e2e
class TestPresentationPipeline:
    """Tests E2E complets: creation -> edition -> verification."""

    async def test_full_pipeline_all_layouts(
        self, generator: PresentationGenerator, editor: PresentationEditor
    ) -> None:
        """Cree une presentation avec 8 layouts, edite, et verifie."""
        # === STEP 1: Create with all 8 layouts ===
        create_params = CreatePresentationParams(
            filename="complete.pptx",
            slides=[
                TitleSlideDef(
                    title="Annual Report 2026",
                    subtitle="Assistance Micro Design",
                    title_style=TextStyle(bold=True, font_size=44),
                    notes="Welcome everyone!",
                ),
                ContentBulletsSlideDef(
                    title="Agenda",
                    bullets=[
                        BulletItem(text="Introduction"),
                        BulletItem(text="Key Results", level=1),
                        BulletItem(text="Next Steps"),
                    ],
                ),
                SectionHeaderSlideDef(
                    title="Part 1: Results",
                    subtitle="Q1-Q4 Performance",
                ),
                ContentWithImageSlideDef(
                    title="Growth Chart",
                    bullets=[BulletItem(text="Revenue up 25%")],
                    image=ImageDef(filename="logo.png"),
                ),
                TwoColumnsSlideDef(
                    title="Comparison",
                    left_bullets=[
                        BulletItem(text="2025"),
                        BulletItem(text="$1M revenue", level=1),
                    ],
                    right_bullets=[
                        BulletItem(text="2026"),
                        BulletItem(text="$1.25M revenue", level=1),
                    ],
                ),
                ImageFullSlideDef(
                    image=ImageDef(filename="diagram.png"),
                    title="Architecture",
                    caption="System overview diagram",
                ),
                ChartSlideDef(
                    title="Revenue by Quarter",
                    chart=PresentationChartDef(
                        chart_type="bar",
                        title="Quarterly Revenue",
                        categories=["Q1", "Q2", "Q3", "Q4"],
                        series=[
                            ChartSeriesDef(name="2025", values=[200, 250, 300, 250]),
                            ChartSeriesDef(name="2026", values=[250, 300, 350, 350]),
                        ],
                    ),
                ),
                ClosingSlideDef(
                    title="Thank You!",
                    subtitle="Questions?",
                    bullets=[BulletItem(text="contact@example.com")],
                ),
            ],
            author="E2E Test",
        )
        create_result = await generator.generate(create_params)

        assert create_result.slides_created == 8
        assert create_result.total_images == 2
        assert create_result.total_charts == 1
        assert create_result.file_size_bytes > 0

        # Verify the created file
        prs = Presentation(create_result.file_path)
        assert len(prs.slides) == 8
        assert prs.core_properties.author == "E2E Test"
        assert prs.slides[0].notes_slide.notes_text_frame.text == "Welcome everyone!"

        # === STEP 2: Edit the presentation ===
        edit_params = EditPresentationParams(
            filename="complete.pptx",
            operations=[
                UpdateTitleOp(
                    slide_index=0,
                    title="Updated Annual Report 2026",
                    style=TextStyle(bold=True, font_color="003366"),
                ),
                UpdateBulletsOp(
                    slide_index=1,
                    bullets=[
                        BulletItem(text="Updated Introduction"),
                        BulletItem(text="Updated Results", level=1),
                        BulletItem(text="Updated Next Steps"),
                        BulletItem(text="New Item", bold=True),
                    ],
                ),
                UpdateNotesOp(slide_index=2, notes="Transition to results section"),
                SetBackgroundOp(slide_index=0, color="F0F0F0"),
                AddSlideOp(
                    slide=ContentBulletsSlideDef(
                        title="Extra Slide",
                        bullets=[BulletItem(text="Added via edit")],
                    ),
                ),
            ],
        )
        edit_result = await editor.edit(edit_params)

        assert edit_result.operations_applied == 5
        assert edit_result.operations_skipped == 0
        assert edit_result.file_size_bytes > 0

        # === STEP 3: Verify the edited file ===
        prs2 = Presentation(edit_result.file_path)
        assert len(prs2.slides) == 9  # 8 original + 1 added

        # Check title was updated
        title_slide = prs2.slides[0]
        texts = [s.text_frame.text for s in title_slide.shapes if s.has_text_frame]
        assert "Updated Annual Report 2026" in texts

        # Check notes on slide 2
        assert prs2.slides[2].notes_slide.notes_text_frame.text == "Transition to results section"

    async def test_create_edit_delete_cycle(
        self, generator: PresentationGenerator, editor: PresentationEditor, tmp_output: Path
    ) -> None:
        """Test cycle complet: create -> add slides -> delete slide."""
        # Create
        create_params = CreatePresentationParams(
            filename="cycle.pptx",
            slides=[
                TitleSlideDef(title="Slide 0"),
                TitleSlideDef(title="Slide 1"),
                TitleSlideDef(title="Slide 2"),
            ],
        )
        await generator.generate(create_params)

        # Delete middle slide
        edit_params = EditPresentationParams(
            filename="cycle.pptx",
            operations=[DeleteSlideOp(slide_index=1)],
        )
        result = await editor.edit(edit_params)
        assert result.operations_applied == 1

        prs = Presentation(str(tmp_output / "cycle.pptx"))
        assert len(prs.slides) == 2
