"""Tests unitaires pour InspectGeneratedFileTool."""

from __future__ import annotations

from pathlib import Path

import pytest
from openpyxl import Workbook
from pptx import Presentation

from src.mcp.tools.inspect_generated_file import InspectGeneratedFileTool


@pytest.fixture
def tool(tmp_path: Path) -> InspectGeneratedFileTool:
    """Tool avec output_path temporaire."""
    tool = InspectGeneratedFileTool()
    tool._output_path = tmp_path
    tool._inspector._output_path = tmp_path
    tool._initialized = True
    return tool


@pytest.fixture
def xlsx_file(tmp_path: Path) -> Path:
    """Cree un fichier Excel de test."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Value"])
    ws.append(["Alice", 100])
    path = tmp_path / "test.xlsx"
    wb.save(path)
    return path


@pytest.fixture
def pptx_file(tmp_path: Path) -> Path:
    """Cree un fichier PowerPoint de test."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Title"
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


class TestToolMetadata:
    """Tests pour les metadonnees du tool."""

    def test_tool_name_and_description(self) -> None:
        assert InspectGeneratedFileTool.name == "inspect_generated_file"
        assert "inspecte" in InspectGeneratedFileTool.description.lower()


class TestInspectXlsx:
    """Tests pour l'inspection de fichiers Excel via le tool."""

    @pytest.mark.asyncio
    async def test_inspect_xlsx_returns_structure(
        self, tool: InspectGeneratedFileTool, xlsx_file: Path
    ) -> None:
        result = await tool.execute({"filename": xlsx_file.name})

        assert result["type"] == "excel"
        assert "sheets" in result
        assert result["filename"] == xlsx_file.name

    @pytest.mark.asyncio
    async def test_editable_with_field_present_excel(
        self, tool: InspectGeneratedFileTool, xlsx_file: Path
    ) -> None:
        result = await tool.execute({"filename": xlsx_file.name})

        assert result["editable_with"] == "edit_excel_document"


class TestInspectPptx:
    """Tests pour l'inspection de fichiers PowerPoint via le tool."""

    @pytest.mark.asyncio
    async def test_inspect_pptx_returns_structure(
        self, tool: InspectGeneratedFileTool, pptx_file: Path
    ) -> None:
        result = await tool.execute({"filename": pptx_file.name})

        assert result["type"] == "presentation"
        assert "slides" in result
        assert result["filename"] == pptx_file.name

    @pytest.mark.asyncio
    async def test_editable_with_field_present_pptx(
        self, tool: InspectGeneratedFileTool, pptx_file: Path
    ) -> None:
        result = await tool.execute({"filename": pptx_file.name})

        assert result["editable_with"] == "edit_presentation"


class TestToolErrors:
    """Tests pour les erreurs du tool."""

    @pytest.mark.asyncio
    async def test_unsupported_extension_returns_error(
        self, tool: InspectGeneratedFileTool, tmp_path: Path
    ) -> None:
        (tmp_path / "file.txt").write_text("nope")

        result = await tool.execute({"filename": "file.txt"})

        assert "error" in result

    @pytest.mark.asyncio
    async def test_file_not_found_returns_error(self, tool: InspectGeneratedFileTool) -> None:
        result = await tool.execute({"filename": "nonexistent.xlsx"})

        assert "error" in result

    @pytest.mark.asyncio
    async def test_path_traversal_blocked(self, tool: InspectGeneratedFileTool) -> None:
        with pytest.raises(Exception, match="ne doit pas contenir"):
            await tool.execute({"filename": "../../../etc/passwd"})

    @pytest.mark.asyncio
    async def test_corrupted_xlsx_returns_error(
        self, tool: InspectGeneratedFileTool, tmp_path: Path
    ) -> None:
        (tmp_path / "corrupt.xlsx").write_text("not real")

        result = await tool.execute({"filename": "corrupt.xlsx"})

        assert "error" in result

    @pytest.mark.asyncio
    async def test_corrupted_pptx_returns_error(
        self, tool: InspectGeneratedFileTool, tmp_path: Path
    ) -> None:
        (tmp_path / "corrupt.pptx").write_text("not real")

        result = await tool.execute({"filename": "corrupt.pptx"})

        assert "error" in result


class TestPptxTitleDetection:
    """Tests pour la coherence de detection titre avec PresentationEditor."""

    @pytest.mark.asyncio
    async def test_pptx_title_detection_matches_editor(self, tmp_path: Path) -> None:
        from src.services.presentation.editor import PresentationEditor

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Titre Coherent"
        path = tmp_path / "coherent.pptx"
        prs.save(str(path))

        # Inspect
        tool = InspectGeneratedFileTool()
        tool._output_path = tmp_path
        tool._inspector._output_path = tmp_path
        tool._initialized = True
        result = await tool.execute({"filename": "coherent.pptx"})

        # Editor detection
        editor = PresentationEditor(output_path=tmp_path)
        prs2 = Presentation(str(path))
        title_shape = editor._find_title_shape(prs2.slides[0])

        assert result["slides"][0]["title"] == title_shape.text_frame.text
