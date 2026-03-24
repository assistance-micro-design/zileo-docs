# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Assistance Micro Design
"""Tests pour PresentationGenerator."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from src.core.exceptions import (
    PresentationGenerationError,
    PresentationImageNotFoundError,
    PresentationOutputTooLargeError,
    PresentationTemplateNotFoundError,
)
from src.models.api import CreatePresentationParams
from src.models.presentation_generation import (
    BulletItem,
    ChartSeriesDef,
    ChartSlideDef,
    ClosingSlideDef,
    ContentBulletsSlideDef,
    ContentWithImageSlideDef,
    ImageDef,
    ImageFullSlideDef,
    PresentationChartDef,
    SectionHeaderSlideDef,
    TextStyle,
    TitleSlideDef,
    TwoColumnsSlideDef,
)
from src.services.presentation.generator import PresentationGenerator


@pytest.fixture
def tmp_output(tmp_path: Path) -> Path:
    """Repertoire de sortie temporaire."""
    output = tmp_path / "output"
    output.mkdir()
    return output


@pytest.fixture
def tmp_images(tmp_path: Path) -> Path:
    """Repertoire d'images temporaire avec une image de test."""
    from tests.conftest import create_minimal_png

    images = tmp_path / "images"
    images.mkdir()
    create_minimal_png(images / "test.png")
    create_minimal_png(images / "photo.jpg")
    return images


@pytest.fixture
def tmp_templates(tmp_path: Path) -> Path:
    """Repertoire de templates temporaire."""
    templates = tmp_path / "templates"
    templates.mkdir()
    # Creer un template minimal
    prs = Presentation()
    prs.save(str(templates / "corporate.pptx"))
    return templates


@pytest.fixture
def generator(tmp_output: Path, tmp_images: Path, tmp_templates: Path) -> PresentationGenerator:
    """Instance PresentationGenerator avec paths temporaires."""
    return PresentationGenerator(
        output_path=tmp_output,
        images_path=tmp_images,
        templates_path=tmp_templates,
    )


class TestSanitizeFilename:
    """Tests pour sanitize_filename."""

    def test_valid_filename(self, generator: PresentationGenerator) -> None:
        assert generator.sanitize_filename("report.pptx") == "report.pptx"

    def test_filename_with_spaces(self, generator: PresentationGenerator) -> None:
        assert generator.sanitize_filename("my report.pptx") == "my report.pptx"

    def test_path_traversal_dotdot(self, generator: PresentationGenerator) -> None:
        with pytest.raises(PresentationGenerationError, match="invalide"):
            generator.sanitize_filename("../evil.pptx")

    def test_path_traversal_slash(self, generator: PresentationGenerator) -> None:
        with pytest.raises(PresentationGenerationError, match="invalide"):
            generator.sanitize_filename("sub/evil.pptx")

    def test_path_traversal_backslash(self, generator: PresentationGenerator) -> None:
        with pytest.raises(PresentationGenerationError, match="invalide"):
            generator.sanitize_filename("sub\\evil.pptx")


class TestEnsureOutputDir:
    """Tests pour ensure_output_dir."""

    def test_creates_directory(self, tmp_path: Path) -> None:
        output = tmp_path / "new_dir" / "sub"
        gen = PresentationGenerator(output_path=output)
        gen.ensure_output_dir()
        assert output.exists()

    def test_existing_directory_ok(self, generator: PresentationGenerator) -> None:
        generator.ensure_output_dir()  # Should not raise


class TestSaveAndVerify:
    """Tests pour save_and_verify."""

    def test_saves_file(self, generator: PresentationGenerator, tmp_output: Path) -> None:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        file_path = tmp_output / "test.pptx"
        size = generator.save_and_verify(prs, file_path, "test.pptx")
        assert file_path.exists()
        assert size > 0

    def test_too_large_file_rejected(self, tmp_output: Path) -> None:
        gen = PresentationGenerator(output_path=tmp_output)
        gen._max_output_size_mb = 0  # 0 MB max -> force error
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        file_path = tmp_output / "big.pptx"
        with pytest.raises(PresentationOutputTooLargeError, match="trop volumineux"):
            gen.save_and_verify(prs, file_path, "big.pptx")
        # File should be deleted
        assert not file_path.exists()


class TestLoadOrCreatePresentation:
    """Tests pour _load_or_create_presentation."""

    def test_no_template(self, generator: PresentationGenerator) -> None:
        prs = generator._load_or_create_presentation(None)
        assert prs is not None
        assert hasattr(prs, "slides")

    def test_valid_template(self, generator: PresentationGenerator) -> None:
        prs = generator._load_or_create_presentation("corporate.pptx")
        assert prs is not None
        assert hasattr(prs, "slides")

    def test_missing_template(self, generator: PresentationGenerator) -> None:
        with pytest.raises(PresentationTemplateNotFoundError, match="introuvable"):
            generator._load_or_create_presentation("nonexistent.pptx")

    def test_template_path_traversal(self, generator: PresentationGenerator) -> None:
        with pytest.raises(PresentationGenerationError, match="invalide"):
            generator._load_or_create_presentation("../evil.pptx")


class TestListTemplates:
    """Tests pour _list_templates."""

    def test_lists_pptx_files(self, generator: PresentationGenerator) -> None:
        templates = generator._list_templates()
        assert "corporate.pptx" in templates

    def test_empty_when_dir_missing(self, tmp_output: Path) -> None:
        gen = PresentationGenerator(
            output_path=tmp_output,
            templates_path=tmp_output / "nonexistent",
        )
        assert gen._list_templates() == []


class TestResolveImagePath:
    """Tests pour _resolve_image_path."""

    def test_valid_image(self, generator: PresentationGenerator) -> None:
        path = generator._resolve_image_path("test.png")
        assert path.exists()

    def test_missing_image(self, generator: PresentationGenerator) -> None:
        with pytest.raises(PresentationImageNotFoundError, match="introuvable"):
            generator._resolve_image_path("nonexistent.png")

    def test_image_path_traversal(self, generator: PresentationGenerator) -> None:
        with pytest.raises(PresentationGenerationError, match="invalide"):
            generator._resolve_image_path("../secret.png")


class TestTitleSlide:
    """Tests pour _create_title_slide."""

    async def test_title_slide_minimal(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[TitleSlideDef(title="Hello World")],
        )
        result = await generator.generate(params)
        assert result.slides_created == 1
        assert result.total_images == 0
        assert result.total_charts == 0

        # Verify the file
        prs = Presentation(result.file_path)
        assert len(prs.slides) == 1

    async def test_title_slide_with_subtitle_and_notes(
        self, generator: PresentationGenerator
    ) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                TitleSlideDef(
                    title="Main Title",
                    subtitle="Subtitle here",
                    notes="Speaker notes",
                    title_style=TextStyle(bold=True, font_size=44),
                )
            ],
        )
        result = await generator.generate(params)
        prs = Presentation(result.file_path)
        slide = prs.slides[0]
        # Check notes
        assert slide.notes_slide.notes_text_frame.text == "Speaker notes"


class TestContentBulletsSlide:
    """Tests pour _create_content_bullets_slide."""

    async def test_bullets_slide(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                ContentBulletsSlideDef(
                    title="Key Points",
                    bullets=[
                        BulletItem(text="First point", bold=True),
                        BulletItem(text="Sub point", level=1),
                        BulletItem(text="Second point"),
                    ],
                )
            ],
        )
        result = await generator.generate(params)
        assert result.slides_created == 1


class TestSectionHeaderSlide:
    """Tests pour _create_section_header_slide."""

    async def test_section_header(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[SectionHeaderSlideDef(title="Section 1", subtitle="Introduction")],
        )
        result = await generator.generate(params)
        assert result.slides_created == 1


class TestClosingSlide:
    """Tests pour _create_closing_slide."""

    async def test_closing_minimal(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[ClosingSlideDef(title="Thank You!")],
        )
        result = await generator.generate(params)
        assert result.slides_created == 1

    async def test_closing_with_bullets(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                ClosingSlideDef(
                    title="Q&A",
                    subtitle="Questions?",
                    bullets=[BulletItem(text="contact@example.com")],
                )
            ],
        )
        result = await generator.generate(params)
        assert result.slides_created == 1


class TestContentWithImageSlide:
    """Tests pour _create_content_with_image_slide."""

    async def test_slide_with_image(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                ContentWithImageSlideDef(
                    title="With Image",
                    bullets=[BulletItem(text="Point")],
                    image=ImageDef(filename="test.png"),
                )
            ],
        )
        result = await generator.generate(params)
        assert result.slides_created == 1
        assert result.total_images == 1

    async def test_slide_with_missing_image(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                ContentWithImageSlideDef(
                    title="With Image",
                    bullets=[BulletItem(text="Point")],
                    image=ImageDef(filename="nonexistent.png"),
                )
            ],
        )
        with pytest.raises(PresentationImageNotFoundError):
            await generator.generate(params)


class TestImageFullSlide:
    """Tests pour _create_image_full_slide."""

    async def test_full_image(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                ImageFullSlideDef(
                    image=ImageDef(filename="test.png"),
                    title="Architecture",
                    caption="Overview diagram",
                )
            ],
        )
        result = await generator.generate(params)
        assert result.total_images == 1


class TestTwoColumnsSlide:
    """Tests pour _create_two_columns_slide."""

    async def test_two_columns(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                TwoColumnsSlideDef(
                    title="Comparison",
                    left_bullets=[BulletItem(text="Left")],
                    right_bullets=[BulletItem(text="Right")],
                )
            ],
        )
        result = await generator.generate(params)
        assert result.slides_created == 1


class TestChartSlide:
    """Tests pour _create_chart_slide."""

    async def test_bar_chart(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                ChartSlideDef(
                    title="Sales",
                    chart=PresentationChartDef(
                        chart_type="bar",
                        categories=["Q1", "Q2", "Q3"],
                        series=[ChartSeriesDef(name="Revenue", values=[100, 150, 200])],
                    ),
                )
            ],
        )
        result = await generator.generate(params)
        assert result.total_charts == 1

    async def test_pie_chart(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="test.pptx",
            slides=[
                ChartSlideDef(
                    title="Distribution",
                    chart=PresentationChartDef(
                        chart_type="pie",
                        categories=["A", "B", "C"],
                        series=[ChartSeriesDef(name="Share", values=[40, 30, 30])],
                    ),
                )
            ],
        )
        result = await generator.generate(params)
        assert result.total_charts == 1

    async def test_all_chart_types(self, generator: PresentationGenerator) -> None:
        for chart_type in ("bar", "line", "pie", "area", "column", "doughnut"):
            params = CreatePresentationParams(
                filename=f"test_{chart_type}.pptx",
                slides=[
                    ChartSlideDef(
                        title=f"Chart {chart_type}",
                        chart=PresentationChartDef(
                            chart_type=chart_type,
                            categories=["A", "B"],
                            series=[ChartSeriesDef(name="S", values=[10, 20])],
                        ),
                    )
                ],
            )
            result = await generator.generate(params)
            assert result.total_charts == 1


class TestMultiSlidePresentation:
    """Tests pour presentations multi-slides."""

    async def test_full_presentation(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="full.pptx",
            slides=[
                TitleSlideDef(title="My Pres", subtitle="By Me"),
                ContentBulletsSlideDef(
                    title="Agenda",
                    bullets=[BulletItem(text="Item 1"), BulletItem(text="Item 2")],
                ),
                SectionHeaderSlideDef(title="Part 1"),
                ContentWithImageSlideDef(
                    title="With Image",
                    bullets=[BulletItem(text="See image")],
                    image=ImageDef(filename="test.png"),
                ),
                TwoColumnsSlideDef(
                    title="Compare",
                    left_bullets=[BulletItem(text="L")],
                    right_bullets=[BulletItem(text="R")],
                ),
                ChartSlideDef(
                    title="Data",
                    chart=PresentationChartDef(
                        chart_type="bar",
                        categories=["A"],
                        series=[ChartSeriesDef(name="S", values=[1])],
                    ),
                ),
                ClosingSlideDef(title="Thanks!"),
            ],
            author="Test Author",
        )
        result = await generator.generate(params)
        assert result.slides_created == 7
        assert result.total_images == 1
        assert result.total_charts == 1

        # Verify the file
        prs = Presentation(result.file_path)
        assert len(prs.slides) == 7
        assert prs.core_properties.author == "Test Author"


class TestOverwrite:
    """Tests pour overwrite detection."""

    async def test_overwrite_flag(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="dup.pptx",
            slides=[TitleSlideDef(title="V1")],
        )
        r1 = await generator.generate(params)
        assert r1.overwritten is False

        r2 = await generator.generate(params)
        assert r2.overwritten is True


class TestTemplateUsage:
    """Tests pour l'utilisation de templates."""

    async def test_with_template(self, generator: PresentationGenerator) -> None:
        params = CreatePresentationParams(
            filename="templated.pptx",
            slides=[TitleSlideDef(title="From Template")],
            template="corporate.pptx",
        )
        result = await generator.generate(params)
        assert result.slides_created == 1


class TestTemplatePlaceholders:
    """Tests pour l'utilisation des placeholders natifs des templates."""

    @pytest.fixture
    def rich_template_dir(self, tmp_path: Path) -> Path:
        """Cree un template avec des layouts riches (placeholders natifs)."""
        tpl_dir = tmp_path / "rich_templates"
        tpl_dir.mkdir()
        # Le template par defaut de python-pptx a des layouts standard
        # avec de vrais placeholders (Title Slide, Title and Content, etc.)
        prs = Presentation()
        prs.save(str(tpl_dir / "rich.pptx"))
        return tpl_dir

    @pytest.fixture
    def rich_generator(
        self, tmp_output: Path, tmp_images: Path, rich_template_dir: Path
    ) -> PresentationGenerator:
        """Generateur avec template riche."""
        return PresentationGenerator(
            output_path=tmp_output,
            images_path=tmp_images,
            templates_path=rich_template_dir,
        )

    async def test_title_slide_uses_native_layout(
        self, rich_generator: PresentationGenerator
    ) -> None:
        """Avec template, title_slide doit utiliser un layout natif, pas Blank."""
        params = CreatePresentationParams(
            filename="native.pptx",
            slides=[TitleSlideDef(title="Native Title", subtitle="Native Sub")],
            template="rich.pptx",
        )
        result = await rich_generator.generate(params)
        prs = Presentation(result.file_path)
        slide = prs.slides[0]
        # Le layout ne doit pas etre Blank
        assert slide.slide_layout.name != "Blank"

    async def test_content_bullets_uses_native_layout(
        self, rich_generator: PresentationGenerator
    ) -> None:
        """Avec template, content_bullets doit utiliser un layout avec BODY."""
        params = CreatePresentationParams(
            filename="native_bullets.pptx",
            slides=[
                ContentBulletsSlideDef(
                    title="Points",
                    bullets=[BulletItem(text="Point 1"), BulletItem(text="Point 2")],
                )
            ],
            template="rich.pptx",
        )
        result = await rich_generator.generate(params)
        prs = Presentation(result.file_path)
        slide = prs.slides[0]
        assert slide.slide_layout.name != "Blank"

    async def test_section_header_uses_native_layout(
        self, rich_generator: PresentationGenerator
    ) -> None:
        params = CreatePresentationParams(
            filename="native_section.pptx",
            slides=[SectionHeaderSlideDef(title="Section", subtitle="Sub")],
            template="rich.pptx",
        )
        result = await rich_generator.generate(params)
        prs = Presentation(result.file_path)
        slide = prs.slides[0]
        assert slide.slide_layout.name != "Blank"

    async def test_two_columns_uses_native_layout(
        self, rich_generator: PresentationGenerator
    ) -> None:
        params = CreatePresentationParams(
            filename="native_cols.pptx",
            slides=[
                TwoColumnsSlideDef(
                    title="Compare",
                    left_bullets=[BulletItem(text="L")],
                    right_bullets=[BulletItem(text="R")],
                )
            ],
            template="rich.pptx",
        )
        result = await rich_generator.generate(params)
        prs = Presentation(result.file_path)
        slide = prs.slides[0]
        assert slide.slide_layout.name != "Blank"

    async def test_template_existing_slides_removed(
        self, rich_generator: PresentationGenerator, tmp_output: Path
    ) -> None:
        """Les slides existants du template doivent etre supprimes."""
        # Le template par defaut de python-pptx n'a pas de slides existants
        # Creons un template avec un slide existant
        tpl_dir = rich_generator._templates_path
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # 1 slide existant
        prs.save(str(tpl_dir / "with_slide.pptx"))

        params = CreatePresentationParams(
            filename="cleaned.pptx",
            slides=[TitleSlideDef(title="Only This")],
            template="with_slide.pptx",
        )
        result = await rich_generator.generate(params)
        prs_out = Presentation(result.file_path)
        # Seul le slide demande doit etre present, pas celui du template
        assert len(prs_out.slides) == 1

    async def test_without_template_uses_blank(
        self, rich_generator: PresentationGenerator
    ) -> None:
        """Sans template, le comportement reste Blank + textbox."""
        params = CreatePresentationParams(
            filename="no_tpl.pptx",
            slides=[TitleSlideDef(title="No Template")],
        )
        result = await rich_generator.generate(params)
        prs = Presentation(result.file_path)
        slide = prs.slides[0]
        assert slide.slide_layout.name == "Blank"

    async def test_full_presentation_with_template(
        self, rich_generator: PresentationGenerator
    ) -> None:
        """Presentation complete avec template riche : aucun crash."""
        params = CreatePresentationParams(
            filename="full_tpl.pptx",
            slides=[
                TitleSlideDef(title="Intro", subtitle="Sub"),
                ContentBulletsSlideDef(
                    title="Points",
                    bullets=[BulletItem(text="A"), BulletItem(text="B", level=1)],
                ),
                SectionHeaderSlideDef(title="Section 1"),
                TwoColumnsSlideDef(
                    title="Compare",
                    left_bullets=[BulletItem(text="L")],
                    right_bullets=[BulletItem(text="R")],
                ),
                ChartSlideDef(
                    title="Data",
                    chart=PresentationChartDef(
                        chart_type="bar",
                        categories=["A"],
                        series=[ChartSeriesDef(name="S", values=[1])],
                    ),
                ),
                ClosingSlideDef(title="Merci"),
            ],
            template="rich.pptx",
        )
        result = await rich_generator.generate(params)
        assert result.slides_created == 6
        prs = Presentation(result.file_path)
        assert len(prs.slides) == 6
