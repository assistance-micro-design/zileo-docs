# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Assistance Micro Design
"""Tests pour PresentationEditor."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from src.core.exceptions import (
    PresentationFileNotFoundError,
    PresentationSlideNotFoundError,
)
from src.models.api import (
    CreatePresentationParams,
    EditPresentationParams,
)
from src.models.presentation_edit import (
    AddSlideOp,
    DeleteSlideOp,
    ReorderSlideOp,
    SetBackgroundOp,
    UpdateBulletsOp,
    UpdateNotesOp,
    UpdateSubtitleOp,
    UpdateTitleOp,
)
from src.models.presentation_generation import (
    BulletItem,
    ContentBulletsSlideDef,
    TextStyle,
    TitleSlideDef,
)
from src.services.presentation.editor import PresentationEditor
from src.services.presentation.generator import PresentationGenerator


@pytest.fixture
def tmp_output(tmp_path: Path) -> Path:
    """Repertoire de sortie temporaire."""
    output = tmp_path / "output"
    output.mkdir()
    return output


@pytest.fixture
def tmp_images(tmp_path: Path) -> Path:
    """Repertoire d'images temporaire."""
    images = tmp_path / "images"
    images.mkdir()
    return images


@pytest.fixture
def tmp_templates(tmp_path: Path) -> Path:
    """Repertoire de templates temporaire."""
    templates = tmp_path / "templates"
    templates.mkdir()
    return templates


@pytest.fixture
def generator(tmp_output: Path, tmp_images: Path, tmp_templates: Path) -> PresentationGenerator:
    """Instance PresentationGenerator."""
    return PresentationGenerator(
        output_path=tmp_output,
        images_path=tmp_images,
        templates_path=tmp_templates,
    )


@pytest.fixture
def editor(tmp_output: Path) -> PresentationEditor:
    """Instance PresentationEditor."""
    return PresentationEditor(output_path=tmp_output)


@pytest.fixture
async def sample_pptx(generator: PresentationGenerator) -> str:
    """Cree un fichier pptx de test et retourne le filename."""
    params = CreatePresentationParams(
        filename="sample.pptx",
        slides=[
            TitleSlideDef(title="Title Slide", subtitle="Subtitle"),
            ContentBulletsSlideDef(
                title="Bullets Slide",
                bullets=[
                    BulletItem(text="First point"),
                    BulletItem(text="Second point"),
                ],
            ),
            TitleSlideDef(title="Third Slide"),
        ],
    )
    await generator.generate(params)
    return "sample.pptx"


class TestFileNotFound:
    """Tests pour fichier introuvable."""

    async def test_missing_file(self, editor: PresentationEditor) -> None:
        params = EditPresentationParams(
            filename="nonexistent.pptx",
            operations=[UpdateTitleOp(slide_index=0, title="T")],
        )
        with pytest.raises(PresentationFileNotFoundError, match="introuvable"):
            await editor.edit(params)


class TestSlideNotFound:
    """Tests pour slide introuvable."""

    async def test_invalid_index(self, editor: PresentationEditor, sample_pptx: str) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[UpdateTitleOp(slide_index=99, title="T")],
        )
        with pytest.raises(PresentationSlideNotFoundError, match="introuvable"):
            await editor.edit(params)


class TestUpdateTitle:
    """Tests pour update_title."""

    async def test_update_title(
        self, editor: PresentationEditor, sample_pptx: str, tmp_output: Path
    ) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[UpdateTitleOp(slide_index=0, title="New Title")],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1
        assert result.operations_skipped == 0

        # Verify
        prs = Presentation(str(tmp_output / sample_pptx))
        slide = prs.slides[0]
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "New Title" in texts

    async def test_update_title_with_style(
        self, editor: PresentationEditor, sample_pptx: str
    ) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[
                UpdateTitleOp(
                    slide_index=0,
                    title="Styled Title",
                    style=TextStyle(bold=True, font_size=48),
                )
            ],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1


class TestUpdateSubtitle:
    """Tests pour update_subtitle."""

    async def test_update_subtitle(self, editor: PresentationEditor, sample_pptx: str) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[UpdateSubtitleOp(slide_index=0, subtitle="New Subtitle")],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1


class TestUpdateBullets:
    """Tests pour update_bullets."""

    async def test_update_bullets(self, editor: PresentationEditor, sample_pptx: str) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[
                UpdateBulletsOp(
                    slide_index=1,
                    bullets=[
                        BulletItem(text="New bullet 1"),
                        BulletItem(text="New bullet 2", level=1),
                    ],
                )
            ],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1


class TestAddSlide:
    """Tests pour add_slide."""

    async def test_add_slide_append(
        self, editor: PresentationEditor, sample_pptx: str, tmp_output: Path
    ) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[AddSlideOp(slide=TitleSlideDef(title="New Slide"))],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1

        prs = Presentation(str(tmp_output / sample_pptx))
        assert len(prs.slides) == 4

    async def test_add_slide_at_index(
        self, editor: PresentationEditor, sample_pptx: str, tmp_output: Path
    ) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[
                AddSlideOp(
                    slide=TitleSlideDef(title="Inserted"),
                    at_index=1,
                )
            ],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1

        prs = Presentation(str(tmp_output / sample_pptx))
        assert len(prs.slides) == 4


class TestDeleteSlide:
    """Tests pour delete_slide."""

    async def test_delete_slide(
        self, editor: PresentationEditor, sample_pptx: str, tmp_output: Path
    ) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[DeleteSlideOp(slide_index=2)],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1

        prs = Presentation(str(tmp_output / sample_pptx))
        assert len(prs.slides) == 2


class TestReorderSlide:
    """Tests pour reorder_slide."""

    async def test_reorder_slide(self, editor: PresentationEditor, sample_pptx: str) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[ReorderSlideOp(from_index=0, to_index=2)],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1

    async def test_reorder_invalid_index(
        self, editor: PresentationEditor, sample_pptx: str
    ) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[ReorderSlideOp(from_index=0, to_index=99)],
        )
        with pytest.raises(PresentationSlideNotFoundError):
            await editor.edit(params)


class TestUpdateNotes:
    """Tests pour update_notes."""

    async def test_update_notes(
        self, editor: PresentationEditor, sample_pptx: str, tmp_output: Path
    ) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[UpdateNotesOp(slide_index=0, notes="Speaker notes here")],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1

        prs = Presentation(str(tmp_output / sample_pptx))
        assert prs.slides[0].notes_slide.notes_text_frame.text == "Speaker notes here"


class TestSetBackground:
    """Tests pour set_background."""

    async def test_set_background(self, editor: PresentationEditor, sample_pptx: str) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[SetBackgroundOp(slide_index=0, color="4472C4")],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 1


class TestMultipleOperations:
    """Tests pour operations multiples."""

    async def test_multiple_ops(self, editor: PresentationEditor, sample_pptx: str) -> None:
        params = EditPresentationParams(
            filename=sample_pptx,
            operations=[
                UpdateTitleOp(slide_index=0, title="Updated Title"),
                UpdateNotesOp(slide_index=0, notes="Notes"),
                SetBackgroundOp(slide_index=1, color="FFFFFF"),
            ],
        )
        result = await editor.edit(params)
        assert result.operations_applied == 3
        assert result.operations_skipped == 0
